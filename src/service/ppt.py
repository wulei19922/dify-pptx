from pptx import Presentation

prs = Presentation()

# 标题幻灯片
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "《中华人民共和国婚姻法》：核心规范"
subtitle = slide.placeholders[1]
subtitle.text = "三方面内容"

# 基本原则幻灯片
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
title_shape = slide.shapes.title
title_shape.text = "基本原则"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "确立婚姻自由（禁止包办、买卖婚姻）"
p = tf.add_paragraph()
p.text = "一夫一妻制度"
p.level = 1
p = tf.add_paragraph()
p.text = "男女平等制度"
p.level = 1
p = tf.add_paragraph()
p.text = "强调保护妇女、儿童及老年人权益"
p.level = 1

# 婚姻关系规范幻灯片
slide = prs.slides.add_slide(bullet_slide_layout)
title_shape = slide.shapes.title
title_shape.text = "婚姻关系规范"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "结婚条件"
p = tf.add_paragraph()
p.text = "男女双方须完全自愿"
p.level = 2
p = tf.add_paragraph()
p.text = "男性需满22周岁、女性满20周岁"
p.level = 2
p = tf.add_paragraph()
p.text = "禁止直系血亲及三代以内旁系血亲结婚"
p.level = 2
p = tf.add_paragraph()
p.text = "禁止行为"
p.level = 1
p = tf.add_paragraph()
p.text = "重婚"
p.level = 2
p = tf.add_paragraph()
p.text = "有配偶者与他人同居"
p.level = 2
p = tf.add_paragraph()
p.text = "家庭暴力等"
p.level = 2
p = tf.add_paragraph()
p.text = "婚姻效力"
p.level = 1
p = tf.add_paragraph()
p.text = "未登记婚姻需补办登记"
p.level = 2
p = tf.add_paragraph()
p.text = "重婚、近亲婚、未达婚龄等情形可导致婚姻无效"
p.level = 2

# 家庭关系幻灯片
slide = prs.slides.add_slide(bullet_slide_layout)
title_shape = slide.shapes.title
title_shape.text = "家庭关系"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "夫妻"
p = tf.add_paragraph()
p.text = "需相互忠实、尊重"
p.level = 2
p = tf.add_paragraph()
p.text = "共同承担子女抚养教育义务"
p.level = 2
p = tf.add_paragraph()
p.text = "家庭成员"
p.level = 1
p = tf.add_paragraph()
p.text = "应敬老爱幼"
p.level = 2
p = tf.add_paragraph()
p.text = "形成平等、文明的相处模式"
p.level = 2

# 总结幻灯片
slide = prs.slides.add_slide(bullet_slide_layout)
title_shape = slide.shapes.title
title_shape.text = "总结"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "婚姻法通过明确权利义务关系，为构建和谐家庭关系提供法律框架。"
tf.paragraphs[0].level = 0

# 保存演示文稿
prs.save('marriage_law.pptx')